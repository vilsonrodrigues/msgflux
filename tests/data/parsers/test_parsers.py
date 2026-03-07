"""Tests for data parsers."""

from io import BytesIO

import pytest

import msgflux as mf


# ---------------------------------------------------------------------------
# Fixtures — all generated in-memory using the target libraries themselves
# ---------------------------------------------------------------------------


@pytest.fixture
def csv_bytes():
    return b"name,age,city\nAlice,30,NYC\nBob,25,LA\n"


@pytest.fixture
def html_bytes():
    return (
        b"<html><head><title>Test Page</title></head>"
        b"<body><h1>Hello</h1><p>World</p>"
        b'<a href="https://example.com">Link</a>'
        b'<img src="/img/logo.png" alt="logo">'
        b"</body></html>"
    )


@pytest.fixture
def email_bytes():
    return (
        b"From: alice@example.com\r\n"
        b"To: bob@example.com\r\n"
        b"Subject: Test Email\r\n"
        b"Date: Mon, 1 Jan 2024 10:00:00 +0000\r\n"
        b"\r\n"
        b"Hello Bob, this is a test email."
    )


@pytest.fixture
def pdf_bytes():
    pypdf = pytest.importorskip("pypdf")
    # Minimal valid PDF with a text content stream
    content = b"BT /F1 12 Tf 100 700 Td (Hello World) Tj ET"
    content_len = len(content)
    raw = (
        b"%PDF-1.4\n"
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R"
        b" /MediaBox [0 0 612 792]"
        b" /Contents 4 0 R"
        b" /Resources << /Font << /F1 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> >> >> >>\nendobj\n"
        b"4 0 obj\n<< /Length " + str(content_len).encode() + b" >>\nstream\n"
        + content + b"\nendstream\nendobj\n"
        b"xref\n0 5\n"
        b"0000000000 65535 f \n"
        b"trailer\n<< /Size 5 /Root 1 0 R >>\nstartxref\n9\n%%EOF\n"
    )
    # Use pypdf's writer for a well-formed PDF instead
    writer = pypdf.PdfWriter()
    writer.add_blank_page(width=612, height=792)
    buf = BytesIO()
    writer.write(buf)
    return buf.getvalue()


@pytest.fixture
def docx_bytes():
    pytest.importorskip("docx")
    from docx import Document

    doc = Document()
    doc.add_heading("Test Document", 0)
    doc.add_paragraph("This is a test paragraph.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Col A"
    table.cell(0, 1).text = "Col B"
    table.cell(1, 0).text = "val 1"
    table.cell(1, 1).text = "val 2"
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_bytes():
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "Subtitle text"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def xlsx_bytes():
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["Name", "Value"])
    ws.append(["Alice", 42])
    ws.append(["Bob", 7])
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------


def test_csv_response_structure(csv_bytes):
    parser = mf.Parser.csv("csv")
    response = parser(csv_bytes)

    assert response.response_type == "csv_parse"
    assert "text" in response.data
    assert "metadata" in response.data


def test_csv_text_is_markdown_table(csv_bytes):
    parser = mf.Parser.csv("csv")
    response = parser(csv_bytes)

    assert "name" in response.data["text"]
    assert "Alice" in response.data["text"]
    assert "|" in response.data["text"]


def test_csv_metadata(csv_bytes):
    parser = mf.Parser.csv("csv")
    response = parser(csv_bytes)

    meta = response.data["metadata"]
    assert meta.num_rows == 2
    assert meta.num_cols == 3
    assert meta.has_header is True


def test_csv_html_format(csv_bytes):
    parser = mf.Parser.csv("csv", table_format="html")
    response = parser(csv_bytes)

    assert "<table>" in response.data["text"]
    assert "<th>" in response.data["text"]


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------


def test_html_response_structure(html_bytes):
    pytest.importorskip("bs4")
    parser = mf.Parser.html("beautifulsoup")
    response = parser(html_bytes)

    assert response.response_type == "html_parse"
    assert "text" in response.data
    assert "links" in response.data
    assert "images" in response.data
    assert "metadata" in response.data


def test_html_text_extraction(html_bytes):
    pytest.importorskip("bs4")
    parser = mf.Parser.html("beautifulsoup")
    response = parser(html_bytes)

    assert "Hello" in response.data["text"]


def test_html_links_extraction(html_bytes):
    pytest.importorskip("bs4")
    parser = mf.Parser.html("beautifulsoup")
    response = parser(html_bytes)

    assert len(response.data["links"]) == 1
    assert response.data["links"][0]["url"] == "https://example.com"


def test_html_images_extraction(html_bytes):
    pytest.importorskip("bs4")
    parser = mf.Parser.html("beautifulsoup")
    response = parser(html_bytes)

    assert len(response.data["images"]) == 1
    assert response.data["images"][0]["alt"] == "logo"


def test_html_metadata(html_bytes):
    pytest.importorskip("bs4")
    parser = mf.Parser.html("beautifulsoup")
    response = parser(html_bytes)

    assert response.data["metadata"].title == "Test Page"
    assert response.data["metadata"].num_links == 1


def test_html_raw_string():
    pytest.importorskip("bs4")
    parser = mf.Parser.html("beautifulsoup")
    response = parser("<html><body><h1>Direct HTML</h1><p>Content</p></body></html>")

    assert "Direct HTML" in response.data["text"]


# ---------------------------------------------------------------------------
# Email
# ---------------------------------------------------------------------------


def test_email_response_structure(email_bytes):
    parser = mf.Parser.email("email")
    response = parser(email_bytes)

    assert response.response_type == "email_parse"
    assert "text" in response.data
    assert "headers" in response.data
    assert "body" in response.data
    assert "metadata" in response.data


def test_email_headers(email_bytes):
    parser = mf.Parser.email("email")
    response = parser(email_bytes)

    headers = response.data["headers"]
    assert "alice@example.com" in headers["from"]
    assert "bob@example.com" in headers["to"]
    assert headers["subject"] == "Test Email"


def test_email_body(email_bytes):
    parser = mf.Parser.email("email")
    response = parser(email_bytes)

    assert "Hello Bob" in response.data["body"]


def test_email_text_is_markdown(email_bytes):
    parser = mf.Parser.email("email")
    response = parser(email_bytes)

    assert "# Email" in response.data["text"]
    assert "**From:**" in response.data["text"]


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def test_pdf_response_structure(pdf_bytes):
    pytest.importorskip("pypdf")
    parser = mf.Parser.pdf("pypdf")
    response = parser(pdf_bytes)

    assert response.response_type == "pdf_parse"
    assert "text" in response.data
    assert "images" in response.data
    assert "metadata" in response.data


def test_pdf_metadata(pdf_bytes):
    pytest.importorskip("pypdf")
    parser = mf.Parser.pdf("pypdf")
    response = parser(pdf_bytes)

    assert response.data["metadata"].num_pages == 1
    assert response.data["metadata"].extraction_mode == "layout"


def test_pdf_images_is_dict(pdf_bytes):
    pytest.importorskip("pypdf")
    parser = mf.Parser.pdf("pypdf")
    response = parser(pdf_bytes)

    assert isinstance(response.data["images"], dict)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def test_docx_response_structure(docx_bytes):
    pytest.importorskip("docx")
    parser = mf.Parser.docx("python_docx")
    response = parser(docx_bytes)

    assert response.response_type == "docx_parse"
    assert "text" in response.data
    assert "images" in response.data
    assert "metadata" in response.data


def test_docx_heading_extraction(docx_bytes):
    pytest.importorskip("docx")
    parser = mf.Parser.docx("python_docx")
    response = parser(docx_bytes)

    assert "Test Document" in response.data["text"]


def test_docx_table_markdown(docx_bytes):
    pytest.importorskip("docx")
    parser = mf.Parser.docx("python_docx")
    response = parser(docx_bytes)

    assert "|" in response.data["text"]
    assert "Col A" in response.data["text"]


def test_docx_table_html(docx_bytes):
    pytest.importorskip("docx")
    parser = mf.Parser.docx("python_docx", table_format="html")
    response = parser(docx_bytes)

    assert "<table>" in response.data["text"]


def test_docx_metadata(docx_bytes):
    pytest.importorskip("docx")
    parser = mf.Parser.docx("python_docx")
    response = parser(docx_bytes)

    assert response.data["metadata"].num_tables == 1


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def test_pptx_response_structure(pptx_bytes):
    pytest.importorskip("pptx")
    parser = mf.Parser.pptx("python_pptx")
    response = parser(pptx_bytes)

    assert response.response_type == "pptx_parse"
    assert "text" in response.data
    assert "images" in response.data
    assert "metadata" in response.data


def test_pptx_title_extraction(pptx_bytes):
    pytest.importorskip("pptx")
    parser = mf.Parser.pptx("python_pptx")
    response = parser(pptx_bytes)

    assert "Test Slide" in response.data["text"]


def test_pptx_slide_marker(pptx_bytes):
    pytest.importorskip("pptx")
    parser = mf.Parser.pptx("python_pptx")
    response = parser(pptx_bytes)

    assert "Slide number: 1" in response.data["text"]


def test_pptx_metadata(pptx_bytes):
    pytest.importorskip("pptx")
    parser = mf.Parser.pptx("python_pptx")
    response = parser(pptx_bytes)

    assert response.data["metadata"].num_slides == 1


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def test_xlsx_response_structure(xlsx_bytes):
    pytest.importorskip("openpyxl")
    parser = mf.Parser.xlsx("openpyxl")
    response = parser(xlsx_bytes)

    assert response.response_type == "xlsx_parse"
    assert "text" in response.data
    assert "images" in response.data
    assert "metadata" in response.data


def test_xlsx_table_markdown(xlsx_bytes):
    pytest.importorskip("openpyxl")
    parser = mf.Parser.xlsx("openpyxl")
    response = parser(xlsx_bytes)

    assert "|" in response.data["text"]
    assert "Name" in response.data["text"]
    assert "Alice" in response.data["text"]


def test_xlsx_table_html(xlsx_bytes):
    pytest.importorskip("openpyxl")
    parser = mf.Parser.xlsx("openpyxl", table_format="html")
    response = parser(xlsx_bytes)

    assert "<table>" in response.data["text"]


def test_xlsx_metadata(xlsx_bytes):
    pytest.importorskip("openpyxl")
    parser = mf.Parser.xlsx("openpyxl")
    response = parser(xlsx_bytes)

    assert response.data["metadata"].num_sheets == 1
