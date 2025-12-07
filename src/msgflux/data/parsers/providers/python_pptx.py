from typing import Any, Dict, List, Optional, Union

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from msgflux.data.parsers.base import BaseParser
from msgflux.data.parsers.registry import register_parser
from msgflux.data.parsers.response import ParserResponse
from msgflux.data.parsers.types import PptxParser
from msgflux.dotdict import dotdict


@register_parser
class PythonPptxPptxParser(BaseParser, PptxParser):
    """Python-PPTX-based PPTX Parser.

    Converts PowerPoint presentations to Markdown format and extracts embedded images.
    Uses python-pptx library for PPTX processing.

    Features:
    - Slide-by-slide extraction
    - Text content including titles and body
    - Image extraction from slides
    - Speaker notes extraction
    - Markdown output format

    Example:
        >>> parser = Parser.pptx("python_pptx")
        >>> response = parser("presentation.pptx")
        >>> print(response.data["text"])
        >>> print(response.data["images"])
    """

    provider = "python_pptx"

    def __init__(
        self,
        *,
        include_notes: Optional[bool] = True,
        encode_images_base64: Optional[bool] = True,
    ):
        """Initialize Python-PPTX parser.

        Args:
            include_notes:
                If True, include speaker notes in the output.
            encode_images_base64:
                If True, encode images as base64 strings.
                If False, keep as raw bytes.
        """
        if Presentation is None:
            raise ImportError(
                "`python-pptx` is not available. Install with `pip install python-pptx`"
            )

        self.include_notes = include_notes
        self.encode_images_base64 = encode_images_base64
        self._initialize()

    def _initialize(self):
        """Initialize parser state."""
        pass

    def __call__(self, data: Union[str, bytes], **_kwargs) -> ParserResponse:
        """Parse a PPTX document.

        Args:
            data:
                PPTX file path, URL, or bytes.
            **kwargs:
                Additional parsing options (currently unused).

        Returns:
            ParserResponse containing:
            - text: Markdown-formatted content
            - images: Dictionary of extracted images
            - metadata: Document metadata (num_slides, etc.)

        Raises:
            FileNotFoundError:
                If file path doesn't exist.
            ValueError:
                If data type is not supported.
        """
        # Validate file type if it's a path
        if isinstance(data, str) and not data.startswith(("http://", "https://")):
            self._validate_file_type(data, ".pptx")

        # Parse the document
        result = self._parse(data)

        # Create response
        response = ParserResponse()
        response.set_response_type("pptx_parse")
        response.add(result)

        return response

    def _parse(self, data: bytes) -> List[Dict[str, Any]]:  # noqa: C901
        """Parse PPTX and extract content.

        Args:
            data:
                PPTX file path or bytes.

        Returns:
            Dictionary with:
            - text: Markdown content
            - images: Dictionary of images
            - metadata: Document metadata
        """
        md_content = ""
        images_dict = {}

        # Load presentation
        if isinstance(data, bytes):
            from io import BytesIO  # noqa: PLC0415

            presentation = Presentation(BytesIO(data))
        else:
            presentation = Presentation(data)

        num_slides = len(presentation.slides)

        # Extract content from each slide
        for idx, slide in enumerate(presentation.slides, 1):
            # Add slide marker
            md_content += f"\n\n<!-- Slide number: {idx} -->\n"

            # Extract title
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                if title_text:
                    md_content += f"\n# {title_text}\n\n"

            # Extract text from all shapes
            for shape in slide.shapes:
                # Skip title (already processed)
                if hasattr(shape, "text") and shape != slide.shapes.title:
                    text = shape.text.strip()
                    if text:
                        md_content += f"{text}\n\n"

                # Extract images
                if hasattr(shape, "image"):
                    img_extension = shape.image.ext
                    img_name = f"slide{idx}_image_{len(images_dict)}.{img_extension}"
                    img_data = shape.image.blob

                    images_dict[img_name] = img_data

                    # Add image reference in markdown
                    md_content += f"![{img_name}]({img_name})\n\n"

            # Extract speaker notes if requested
            if self.include_notes and slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    md_content += f"\n**Speaker Notes:**\n{notes_text}\n\n"

        # Prepare images (optionally encode to base64)
        images_dict = self._prepare_images_dict(
            images_dict, encode_base64=self.encode_images_base64
        )

        # Prepare metadata
        metadata = dotdict(
            {
                "num_slides": num_slides,
                "include_notes": self.include_notes,
            }
        )

        # Try to extract presentation metadata if available
        if hasattr(presentation, "core_properties"):
            props = presentation.core_properties
            pptx_metadata = {}
            if props.title:
                pptx_metadata["title"] = props.title
            if props.author:
                pptx_metadata["author"] = props.author
            if props.subject:
                pptx_metadata["subject"] = props.subject
            if props.created:
                pptx_metadata["created"] = str(props.created)
            if props.modified:
                pptx_metadata["modified"] = str(props.modified)
            if pptx_metadata:
                metadata.pptx_info = pptx_metadata

        return {
            "text": md_content.strip(),
            "images": images_dict,
            "metadata": metadata,
        }

    async def acall(self, data: Union[str, bytes], **_kwargs) -> ParserResponse:
        """Async version of __call__. Parse a PPTX document asynchronously.

        Args:
            data:
                PPTX file path, URL, or bytes.
            **kwargs:
                Additional parsing options (currently unused).

        Returns:
            ParserResponse containing:
            - text: Markdown-formatted content
            - images: Dictionary of extracted images
            - metadata: Document metadata (num_slides, etc.)

        Raises:
            FileNotFoundError:
                If file path doesn't exist.
            ValueError:
                If data type is not supported.
        """
        # Validate file type if it's a path
        if isinstance(data, str) and not data.startswith(("http://", "https://")):
            self._validate_file_type(data, ".pptx")

        # Load file asynchronously if needed
        if isinstance(data, str):
            data = await self._aload_file(data)

        # Parse the document
        result = self._parse(data)

        # Create response
        response = ParserResponse()
        response.set_response_type("pptx_parse")
        response.add(result)

        return response
